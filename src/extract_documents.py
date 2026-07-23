from __future__ import annotations

import argparse
import io
import json
import logging
import os
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from azure.ai.documentintelligence import DocumentIntelligenceClient
from azure.core.credentials import AzureKeyCredential
from azure.identity import DefaultAzureCredential
from azure.storage.blob import BlobServiceClient, ContentSettings
from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from common import configure_logging, processing_allowed, safe_filename


def blob_service() -> BlobServiceClient:
    return BlobServiceClient(
        account_url=os.environ["AZURE_STORAGE_ACCOUNT_URL"],
        credential=DefaultAzureCredential(),
    )


def document_client() -> DocumentIntelligenceClient | None:
    endpoint = os.getenv("AZURE_DOCUMENT_INTELLIGENCE_ENDPOINT")
    if not endpoint:
        return None
    key = os.getenv("AZURE_DOCUMENT_INTELLIGENCE_KEY")
    credential = AzureKeyCredential(key) if key else DefaultAzureCredential()
    return DocumentIntelligenceClient(endpoint=endpoint, credential=credential)


def extract_docx(data: bytes) -> tuple[str, list[dict[str, Any]]]:
    doc = Document(io.BytesIO(data))
    blocks = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    tables = []
    for index, table in enumerate(doc.tables, start=1):
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        tables.append({"table": index, "rows": rows})
    return "\n\n".join(blocks), tables


def extract_pptx(data: bytes) -> tuple[str, list[dict[str, Any]]]:
    deck = Presentation(io.BytesIO(data))
    slides = []
    blocks = []
    for index, slide in enumerate(deck.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        slides.append({"slide": index, "text": texts})
        blocks.append(f"## Slide {index}\n" + "\n".join(texts))
    return "\n\n".join(blocks), slides


def extract_xlsx(data: bytes) -> tuple[str, list[dict[str, Any]]]:
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sheets = []
    blocks = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(values):
                rows.append(values)
        sheets.append({"sheet": ws.title, "rows": rows})
        blocks.append(f"## Sheet: {ws.title}\n" + "\n".join(" | ".join(row) for row in rows))
    return "\n\n".join(blocks), sheets


def extract_pdf_fallback(data: bytes) -> tuple[str, list[dict[str, Any]]]:
    reader = PdfReader(io.BytesIO(data))
    pages = []
    blocks = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        pages.append({"page": index, "text": text})
        blocks.append(f"## Page {index}\n{text}")
    return "\n\n".join(blocks), pages


def extract_with_document_intelligence(client: DocumentIntelligenceClient, data: bytes) -> tuple[str, list[dict[str, Any]]]:
    poller = client.begin_analyze_document("prebuilt-layout", body=io.BytesIO(data))
    result = poller.result()
    tables = []
    for t_index, table in enumerate(result.tables or [], start=1):
        cells = []
        for cell in table.cells:
            cells.append({
                "row": cell.row_index,
                "column": cell.column_index,
                "content": cell.content,
                "row_span": cell.row_span,
                "column_span": cell.column_span,
            })
        tables.append({"table": t_index, "row_count": table.row_count, "column_count": table.column_count, "cells": cells})
    return result.content or "", tables


def extract_file(name: str, content_type: str, data: bytes, di_client: DocumentIntelligenceClient | None):
    ext = Path(name).suffix.lower()
    if ext == ".docx":
        return extract_docx(data)
    if ext == ".pptx":
        return extract_pptx(data)
    if ext == ".xlsx":
        return extract_xlsx(data)
    if ext in {".csv", ".txt", ".md"}:
        return data.decode("utf-8", errors="replace"), []
    if ext == ".pdf":
        if di_client:
            return extract_with_document_intelligence(di_client, data)
        return extract_pdf_fallback(data)
    if ext in {".jpg", ".jpeg", ".png"} and di_client:
        return extract_with_document_intelligence(di_client, data)
    raise ValueError(f"Unsupported extraction type: {ext or content_type}")


def main() -> int:
    parser = argparse.ArgumentParser(description="Extract text and tables from raw Azure Blob documents.")
    parser.add_argument("--prefix", default="raw/")
    parser.add_argument("--verbose", action="store_true")
    args = parser.parse_args()
    configure_logging(args.verbose)

    service = blob_service()
    container = service.get_container_client(os.getenv("AZURE_STORAGE_CONTAINER", "aari-knowledge-base"))
    di_client = document_client()
    run_stamp = datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%SZ")

    processed = 0
    errors = 0
    for blob in container.list_blobs(name_starts_with=args.prefix, include=["metadata"]):
        metadata = blob.metadata or {}
        if not processing_allowed(metadata, blob.name):
            logging.info("Policy skip (not AI-processable): %s", blob.name)
            continue
        try:
            client = container.get_blob_client(blob.name)
            data = client.download_blob().readall()
            content_type = getattr(blob, "content_settings", None)
            mime = getattr(content_type, "content_type", "") if content_type else ""
            text, structures = extract_file(blob.name, mime, data, di_client)
            record = {
                "source_blob": blob.name,
                "source_metadata": metadata,
                "extracted_at": datetime.now(timezone.utc).isoformat(),
                "content": text,
                "structures": structures,
            }
            category = metadata.get("category", "uncategorized")
            base = safe_filename(Path(blob.name).name)
            processed_json = f"processed/{category}/{base}.json"
            processed_md = f"processed/{category}/{base}.md"
            container.upload_blob(
                processed_json,
                json.dumps(record, ensure_ascii=False, indent=2).encode("utf-8"),
                overwrite=True,
                metadata={**metadata, "source_blob": blob.name[:1500]},
                content_settings=ContentSettings(content_type="application/json"),
            )
            container.upload_blob(
                processed_md,
                text.encode("utf-8"),
                overwrite=True,
                metadata={**metadata, "source_blob": blob.name[:1500]},
                content_settings=ContentSettings(content_type="text/markdown"),
            )
            processed += 1
            logging.info("Extracted %s", blob.name)
        except Exception as exc:
            errors += 1
            logging.exception("Failed extracting %s: %s", blob.name, exc)

    summary = {"run": run_stamp, "processed": processed, "errors": errors}
    container.upload_blob(
        f"manifests/{run_stamp}/extraction_summary.json",
        json.dumps(summary, indent=2).encode("utf-8"),
        overwrite=True,
        content_settings=ContentSettings(content_type="application/json"),
    )
    logging.info("Extraction complete: %s", summary)
    return 1 if errors else 0


if __name__ == "__main__":
    raise SystemExit(main())
