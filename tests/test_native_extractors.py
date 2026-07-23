from pathlib import Path
import io
import sys

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from extract_documents import extract_docx, extract_pptx, extract_xlsx


def test_docx_extraction():
    doc = Document()
    doc.add_paragraph("AARI curriculum")
    table = doc.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Week"
    table.cell(0, 1).text = "Topic"
    stream = io.BytesIO()
    doc.save(stream)
    text, tables = extract_docx(stream.getvalue())
    assert "AARI curriculum" in text
    assert tables[0]["rows"][0] == ["Week", "Topic"]


def test_xlsx_extraction():
    wb = Workbook()
    ws = wb.active
    ws.title = "Budget"
    ws.append(["Item", "Cost"])
    ws.append(["Server", 1000])
    stream = io.BytesIO()
    wb.save(stream)
    text, sheets = extract_xlsx(stream.getvalue())
    assert "Server" in text
    assert sheets[0]["sheet"] == "Budget"


def test_pptx_extraction():
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = "AARI"
    slide.placeholders[1].text = "Operators, not observers"
    stream = io.BytesIO()
    deck.save(stream)
    text, slides = extract_pptx(stream.getvalue())
    assert "Operators, not observers" in text
    assert slides[0]["slide"] == 1
